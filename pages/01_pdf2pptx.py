from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import os

def pdf_to_pptx(pdf_path, output_path, poppler_path=None):
    """
    Converts a PDF file to a PPTX file by converting each page into an image.

    :param pdf_path: Path to the input PDF file
    :param output_path: Path to the output PPTX file
    :param poppler_path: Path to the poppler binaries (for Windows users)
    """
    # Convert PDF pages to images
    images = convert_from_path(pdf_path, poppler_path=poppler_path)
    
    # Create a PowerPoint presentation
    presentation = Presentation()

    for image in images:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # Blank slide layout
        image_path = "temp_image.png"
        image.save(image_path, "PNG")

        left = Inches(0)
        top = Inches(0)
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height

        slide.shapes.add_picture(image_path, left, top, slide_width, slide_height)

        # Remove temporary image file
        os.remove(image_path)

    # Save the PowerPoint file
    presentation.save(output_path)

# Streamlit App
st.title("PDF to PPTX Converter")

uploaded_file = st.file_uploader("Upload a PDF file", type="pdf")

if uploaded_file is not None:
    # Save the uploaded file temporarily
    pdf_path = "uploaded_file.pdf"
    with open(pdf_path, "wb") as f:
        f.write(uploaded_file.read())

    # Convert the PDF to PPTX
    output_path = "output_presentation.pptx"
    poppler_path = None
    if os.name == "nt":  # Windows
        poppler_path = r"C:\\path\\to\\poppler\\bin"  # Adjust this to your Poppler path
    
    pdf_to_pptx(pdf_path, output_path, poppler_path)

    # Provide the PPTX file for download
    with open(output_path, "rb") as f:
        st.download_button(
            label="Download PPTX",
            data=f,
            file_name="converted_presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    # Clean up temporary files
    os.remove(pdf_path)
    os.remove(output_path)
